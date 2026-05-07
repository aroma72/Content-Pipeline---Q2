"""
Systems Evaluations Instructor Slides - CLEAN VERSION (No Images)
Professional minimal design with excellent typography and formatting
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Professional colors
PRIMARY = RGBColor(31, 78, 121)      # Dark blue
ACCENT = RGBColor(0, 176, 176)       # Teal
LIGHT_BG = RGBColor(240, 248, 255)   # Light blue
TEXT_DARK = RGBColor(50, 50, 50)     # Dark gray
TEXT_MED = RGBColor(100, 100, 100)   # Medium gray
WHITE = RGBColor(255, 255, 255)

def add_header(slide, title):
    """Add professional header"""
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.85))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY
    header.line.color.rgb = PRIMARY

    tf = header.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.15)

def add_accent_line(slide):
    """Add accent line under header"""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.85), Inches(10), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT

def add_text_box(slide, left, top, width, height, text, size=16, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT):
    """Add formatted text box"""
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.line_spacing = 1.4
    return box

def add_content_box(slide, left, top, width, height, title, content):
    """Add content box with title and content"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BG
    box.line.color.rgb = ACCENT
    box.line.width = Pt(2)

    # Title
    tf_title = box.text_frame
    tf_title.clear()
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.size = Pt(16)
    p_title.font.bold = True
    p_title.font.color.rgb = PRIMARY
    tf_title.margin_left = Inches(0.2)
    tf_title.margin_top = Inches(0.15)

    # Content below
    add_text_box(slide, left + Inches(0.2), top + Inches(0.4), width - Inches(0.4), height - Inches(0.45),
                content, size=13, color=TEXT_DARK)

def create_presentation():
    """Create the presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # SLIDE 1: Title
    print("📄 Slide 1: Title")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = PRIMARY

    add_text_box(s, Inches(0.5), Inches(2.5), Inches(9), Inches(1.5),
                "Systems Evaluations", size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(4.2), Inches(9), Inches(1),
                "Instructor Guide", size=40, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(5.5), Inches(9), Inches(0.8),
                "Professional Quality Control for AI Agents", size=18, color=RGBColor(200, 200, 200), align=PP_ALIGN.CENTER)

    # SLIDE 2: What is Evaluation?
    print("📄 Slide 2: What is Evaluation?")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    add_header(s, "What is Evaluation?")
    add_accent_line(s)

    add_text_box(s, Inches(0.4), Inches(1.3), Inches(4.8), Inches(0.5),
                "Quality control for agents.", size=20, bold=True)

    add_text_box(s, Inches(0.4), Inches(2), Inches(4.8), Inches(5),
                "You ask:\n\n• Does it work?\n• Does it meet standards?\n• Where does it fail?\n• How do you improve it?",
                size=16, color=TEXT_DARK)

    # Right content box
    add_content_box(s, Inches(5.4), Inches(1.3), Inches(4.2), Inches(5.8),
                "Unlike Testing",
                "Testing checks if code works.\n\nEvaluation measures whether goals are achieved — quality, helpfulness, and reliability in the real world.")

    # SLIDE 3: Types of Evaluations
    print("📄 Slide 3: Types of Evaluations")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    add_header(s, "Types of Evaluations")
    add_accent_line(s)

    # Three columns
    cols = [
        ("Automated\nEvaluation", "• Speed and scale\n• Defined metrics\n• Repeatable\n• Consistent"),
        ("Manual\nReview", "• Nuance and context\n• Human judgment\n• Quality gates\n• Edge cases"),
        ("Hybrid\nApproach", "• Automated first\n• Human spot-check\n• Best of both\n• Cost effective")
    ]

    x_start = 0.5
    for i, (title, content) in enumerate(cols):
        left = Inches(x_start + (i * 3))
        add_content_box(s, left, Inches(1.3), Inches(2.8), Inches(5.8), title, content)

    # SLIDE 4: Evaluation Framework
    print("📄 Slide 4: Evaluation Framework")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    add_header(s, "Evaluation Framework")
    add_accent_line(s)

    steps = [
        ("1. Define\nSuccess", "What does good look like?"),
        ("2. Set\nMetrics", "How do you measure it?"),
        ("3. Create\nRubrics", "What's the scoring?"),
        ("4. Run\nTests", "Test on samples"),
        ("5. Collect\nFeedback", "Iterate & refine")
    ]

    y = 1.3
    for title, content in steps:
        add_content_box(s, Inches(0.4), Inches(y), Inches(9.2), Inches(0.95), title, content)
        y += 1.05

    # SLIDE 5: Rubrics & Scoring
    print("📄 Slide 5: Creating Effective Rubrics")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    add_header(s, "Creating Effective Rubrics")
    add_accent_line(s)

    rubric_items = [
        ("Criteria", "Clear, measurable dimensions of quality"),
        ("Levels", "Scale (1-4: Exemplary, Proficient, Developing, Needs Work)"),
        ("Descriptors", "Specific examples of what each level looks like"),
        ("Evidence", "Real sample outputs showing each level")
    ]

    y = 1.3
    for label, desc in rubric_items:
        add_content_box(s, Inches(0.4), Inches(y), Inches(9.2), Inches(1.2), label, desc)
        y += 1.3

    # SLIDE 6: Facilitation Tips
    print("📄 Slide 6: Facilitation Tips")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    add_header(s, "Facilitation Tips")
    add_accent_line(s)

    tips = [
        ("Be Specific", "Use concrete examples, not vague feedback"),
        ("Focus on Criteria", "Anchor judgments to your defined metrics"),
        ("Calibrate Raters", "Discuss edge cases to ensure consistency"),
        ("Document Decisions", "Log your rationale for audit trail"),
        ("Iterate", "Use feedback to refine your rubrics")
    ]

    y = 1.3
    for tip, detail in tips:
        add_content_box(s, Inches(0.4), Inches(y), Inches(9.2), Inches(0.95), tip, detail)
        y += 1.0

    # SLIDE 7: Common Pitfalls
    print("📄 Slide 7: Common Pitfalls to Avoid")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    add_header(s, "Common Pitfalls to Avoid")
    add_accent_line(s)

    pitfalls = [
        ("Unmeasurable Criteria", "Avoid vague terms like 'good' or 'bad'"),
        ("Scope Creep", "Test one thing at a time, not everything"),
        ("Insufficient Samples", "Need enough data for statistical confidence"),
        ("Rater Bias", "Use structured rubrics, not gut feeling"),
        ("No Documentation", "Always log why a decision was made")
    ]

    y = 1.3
    for pitfall, solution in pitfalls:
        # Warning box (red)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(y), Inches(9.2), Inches(0.95))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 245, 245)
        box.line.color.rgb = RGBColor(220, 53, 69)
        box.line.width = Pt(2)

        add_text_box(s, Inches(0.6), Inches(y + 0.08), Inches(8.8), Inches(0.3),
                    f"❌ {pitfall}", size=14, bold=True, color=RGBColor(220, 53, 69))
        add_text_box(s, Inches(0.8), Inches(y + 0.45), Inches(8.4), Inches(0.4),
                    f"✓ {solution}", size=12, color=TEXT_DARK)
        y += 1.0

    # SLIDE 8: Getting Started
    print("📄 Slide 8: Getting Started")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    add_header(s, "Getting Started")
    add_accent_line(s)

    steps_start = [
        ("Step 1", "Define your evaluation goals\nWhat are you evaluating?"),
        ("Step 2", "Design your rubric\nBuild your measurement tool"),
        ("Step 3", "Pilot test\nRun on sample outputs"),
        ("Step 4", "Refine\nAdjust based on results"),
        ("Step 5", "Scale up\nApply to full dataset")
    ]

    y = 1.3
    for step, desc in steps_start:
        add_content_box(s, Inches(0.4), Inches(y), Inches(9.2), Inches(1.05), step, desc)
        y += 1.1

    # SLIDE 9: Key Takeaways
    print("📄 Slide 9: Key Takeaways")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    add_header(s, "Key Takeaways")
    add_accent_line(s)

    takeaways = [
        "Evaluation measures real-world impact, not just code correctness",
        "Structured rubrics eliminate bias and ensure consistency",
        "Hybrid approaches (automated + manual) work best",
        "Document your decisions for reproducibility",
        "Iterate based on pilot results before scaling"
    ]

    y = 1.5
    for i, takeaway in enumerate(takeaways, 1):
        add_text_box(s, Inches(1), Inches(y), Inches(8), Inches(0.9),
                    f"{i}. {takeaway}", size=18, bold=True, color=PRIMARY)
        y += 1.05

    # SLIDE 10: Q&A
    print("📄 Slide 10: Q&A")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = PRIMARY

    add_text_box(s, Inches(0.5), Inches(2.8), Inches(9), Inches(1.5),
                "Questions?", size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(4.5), Inches(9), Inches(1),
                "aroma.tahir@taleemabad.com", size=24, color=ACCENT, align=PP_ALIGN.CENTER)

    return prs

# Create presentation
print("\n" + "="*70)
print("📊 CREATING SYSTEMS EVALUATIONS INSTRUCTOR SLIDES")
print("   Clean Design • No Images • Professional Typography")
print("="*70 + "\n")

prs = create_presentation()

from pathlib import Path
output_dir = Path("weekly_artifacts/week-20-2026")
output_dir.mkdir(parents=True, exist_ok=True)

output_path = output_dir / "Systems_Evaluations_Instructor_Slides_Final.pptx"
prs.save(output_path)

print("="*70)
print("✅ COMPLETE!")
print(f"\n   File: Systems_Evaluations_Instructor_Slides_Final.pptx")
print("\n   Content:")
print("   • 10 professional slides")
print("   • Clean minimal design")
print("   • No images (text + formatting focused)")
print("   • Professional color scheme (navy + teal)")
print("   • Ready for instructor delivery")
print("="*70 + "\n")

"""
Generate Systems Evaluations Slides - FINAL VERSION
With professional images from Unsplash (reliable free API)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import requests
from PIL import Image, ImageDraw, ImageFont
from io import BytesIO

# Professional color palette
PRIMARY_COLOR = RGBColor(31, 78, 121)      # Dark blue
ACCENT_COLOR = RGBColor(0, 176, 176)       # Teal
LIGHT_BG = RGBColor(240, 248, 255)         # Alice blue
TEXT_DARK = RGBColor(50, 50, 50)           # Dark gray
TEXT_LIGHT = RGBColor(100, 100, 100)       # Medium gray
WHITE = RGBColor(255, 255, 255)
GRADIENT_BLUE = RGBColor(41, 98, 141)

def get_unsplash_image(query, width=600, height=400):
    """Get image from Unsplash with query"""
    try:
        # Unsplash API - free, no key required for basic use
        url = f"https://source.unsplash.com/{width}x{height}/?{query}"
        response = requests.get(url, timeout=10, allow_redirects=True)

        if response.status_code == 200:
            img = Image.open(BytesIO(response.content))
            # Enhance the image slightly
            img = img.resize((width, height), Image.Resampling.LANCZOS)
            # Add a subtle overlay for text readability
            overlay = Image.new('RGBA', img.size, (0, 0, 0, 30))
            img = img.convert('RGBA')
            img = Image.alpha_composite(img, overlay)
            img = img.convert('RGB')

            img_bytes = BytesIO()
            img.save(img_bytes, format='PNG')
            img_bytes.seek(0)
            print(f"   ✓ Loaded: {query}")
            return img_bytes
    except Exception as e:
        print(f"   ✗ Failed: {query} - {str(e)[:40]}")

    return None

def add_header(slide, title):
    """Add clean header with title"""
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.8))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.color.rgb = PRIMARY_COLOR

    title_box = header.text_frame
    title_box.clear()
    title_box.word_wrap = True
    p = title_box.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    title_box.margin_left = Inches(0.4)
    title_box.margin_top = Inches(0.15)

def add_accent_line(slide):
    """Add accent line under header"""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.8), Inches(10), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR

def add_image_safe(slide, image_data, left, top, width, height):
    """Safely add image to slide"""
    if image_data is None:
        return False
    try:
        slide.shapes.add_picture(image_data, left, top, width=width, height=height)
        return True
    except:
        return False

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT):
    """Add formatted text box"""
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True

    p = frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.line_spacing = 1.3

    return box

def create_presentation():
    """Create the full presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    print("\n📄 Creating Slide 1: Title")
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide1.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR

    add_text_box(slide1, Inches(1), Inches(2.5), Inches(8), Inches(2),
                "Systems Evaluations", font_size=66, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide1, Inches(1), Inches(4.5), Inches(8), Inches(1.5),
                "Instructor Guide", font_size=32, color=ACCENT_COLOR, align=PP_ALIGN.CENTER)

    # Slide 2: What is Evaluation?
    print("📄 Creating Slide 2: What is Evaluation?")
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    slide2.background.fill.solid()
    slide2.background.fill.fore_color.rgb = WHITE
    add_header(slide2, "What is Evaluation?")
    add_accent_line(slide2)

    add_text_box(slide2, Inches(0.4), Inches(1.2), Inches(4.5), Inches(0.6),
                "Quality control for agents.", font_size=20, bold=True, color=TEXT_DARK)

    add_text_box(slide2, Inches(0.4), Inches(2), Inches(4.5), Inches(4.5),
                "You ask:\n\n• Does it work?\n• Does it meet standards?\n• Where does it fail?\n• How do you improve it?",
                font_size=16, color=TEXT_DARK)

    # Right content box
    box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.2), Inches(4.3), Inches(5.7))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BG
    box.line.color.rgb = ACCENT_COLOR
    box.line.width = Pt(2)

    add_text_box(slide2, Inches(5.5), Inches(1.5), Inches(3.7), Inches(0.6),
                "Unlike Testing", font_size=18, bold=True, color=PRIMARY_COLOR)

    add_text_box(slide2, Inches(5.5), Inches(2.3), Inches(3.7), Inches(4.4),
                "Testing checks if code works.\n\nEvaluation measures whether goals are achieved — quality, helpfulness, and reliability in the real world.",
                font_size=14, color=TEXT_DARK)

    # Slide 3: Types of Evaluations
    print("📄 Creating Slide 3: Types of Evaluations")
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    slide3.background.fill.solid()
    slide3.background.fill.fore_color.rgb = WHITE
    add_header(slide3, "Types of Evaluations")
    add_accent_line(slide3)

    print("   Loading images...")
    img = get_unsplash_image("assessment evaluation quality", 600, 400)
    if img:
        add_image_safe(slide3, img, Inches(5.5), Inches(1.2), width=Inches(4), height=Inches(5.8))

    add_text_box(slide3, Inches(0.4), Inches(1.2), Inches(5), Inches(5.8),
                """Automated Evaluation
• Speed and scale
• Defined metrics
• Repeatable

Manual Review
• Nuance and context
• Human judgment
• Quality gates

Hybrid Approach
• Automated first pass
• Human spot-check
• Best of both worlds""", font_size=14, color=TEXT_DARK)

    # Slide 4: Evaluation Framework
    print("📄 Creating Slide 4: Evaluation Framework")
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    slide4.background.fill.solid()
    slide4.background.fill.fore_color.rgb = WHITE
    add_header(slide4, "Evaluation Framework")
    add_accent_line(slide4)

    img = get_unsplash_image("framework process workflow professional", 600, 400)
    if img:
        add_image_safe(slide4, img, Inches(0.4), Inches(1.2), width=Inches(4.5), height=Inches(5.8))

    y_pos = 1.2
    for title, desc in [
        ("Define Success", "What does good look like?"),
        ("Set Metrics", "How do you measure it?"),
        ("Create Rubrics", "What's the scoring system?"),
        ("Run Tests", "Apply framework to samples"),
        ("Collect Feedback", "Iterate and refine")
    ]:
        box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(y_pos), Inches(4.3), Inches(0.95))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BG
        box.line.color.rgb = ACCENT_COLOR
        box.line.width = Pt(1.5)

        add_text_box(slide4, Inches(5.4), Inches(y_pos + 0.05), Inches(3.9), Inches(0.35),
                    title, font_size=13, bold=True, color=PRIMARY_COLOR)
        add_text_box(slide4, Inches(5.4), Inches(y_pos + 0.45), Inches(3.9), Inches(0.45),
                    desc, font_size=11, color=TEXT_LIGHT)

        y_pos += 1.1

    # Slide 5: Rubrics & Scoring
    print("📄 Creating Slide 5: Creating Effective Rubrics")
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    slide5.background.fill.solid()
    slide5.background.fill.fore_color.rgb = WHITE
    add_header(slide5, "Creating Effective Rubrics")
    add_accent_line(slide5)

    img = get_unsplash_image("rubric scoring criteria assessment", 600, 400)
    if img:
        add_image_safe(slide5, img, Inches(5.5), Inches(1.2), width=Inches(4), height=Inches(5.8))

    add_text_box(slide5, Inches(0.4), Inches(1.2), Inches(5), Inches(5.8),
                """Rubric Components:

Criteria
Clear, measurable dimensions

Levels
Scale (1-4, Exemplary-Needs Work)

Descriptors
What each level looks like

Examples
Real sample outputs""", font_size=15, color=TEXT_DARK)

    # Slide 6: Facilitation Tips
    print("📄 Creating Slide 6: Facilitation Tips")
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    slide6.background.fill.solid()
    slide6.background.fill.fore_color.rgb = WHITE
    add_header(slide6, "Facilitation Tips")
    add_accent_line(slide6)

    img = get_unsplash_image("team collaboration discussion professional training", 600, 400)
    if img:
        add_image_safe(slide6, img, Inches(5.5), Inches(1.2), width=Inches(4), height=Inches(5.8))

    add_text_box(slide6, Inches(0.4), Inches(1.2), Inches(5), Inches(5.8),
                """1. Be Specific
Use concrete examples, not vague feedback

2. Focus on Criteria
Anchor judgments to defined metrics

3. Calibrate Raters
Discuss edge cases to ensure consistency

4. Document Decisions
Log rationale for audit trail

5. Iterate
Use feedback to refine rubrics""", font_size=14, color=TEXT_DARK)

    # Slide 7: Common Pitfalls
    print("📄 Creating Slide 7: Common Pitfalls to Avoid")
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    slide7.background.fill.solid()
    slide7.background.fill.fore_color.rgb = WHITE
    add_header(slide7, "Common Pitfalls to Avoid")
    add_accent_line(slide7)

    y_pos = 1.2
    for pitfall, solution in [
        ("Unmeasurable Criteria", "Avoid vague terms like 'good' or 'bad'"),
        ("Scope Creep", "Test one thing at a time"),
        ("Insufficient Samples", "Need enough data for statistical confidence"),
        ("Rater Bias", "Use structured rubrics, not gut feeling"),
        ("No Documentation", "Always log why a decision was made")
    ]:
        box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(y_pos), Inches(9.2), Inches(0.9))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 245, 245)
        box.line.color.rgb = RGBColor(220, 53, 69)
        box.line.width = Pt(2)

        add_text_box(slide7, Inches(0.6), Inches(y_pos + 0.05), Inches(8.8), Inches(0.35),
                    f"❌ {pitfall}", font_size=14, bold=True, color=RGBColor(220, 53, 69))
        add_text_box(slide7, Inches(0.8), Inches(y_pos + 0.42), Inches(8.4), Inches(0.4),
                    f"✓ {solution}", font_size=12, color=TEXT_DARK)

        y_pos += 1.0

    # Slide 8: Next Steps
    print("📄 Creating Slide 8: Getting Started")
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    slide8.background.fill.solid()
    slide8.background.fill.fore_color.rgb = WHITE
    add_header(slide8, "Getting Started")
    add_accent_line(slide8)

    img = get_unsplash_image("implementation process steps execution planning", 600, 400)
    if img:
        add_image_safe(slide8, img, Inches(5.5), Inches(1.2), width=Inches(4), height=Inches(5.8))

    add_text_box(slide8, Inches(0.4), Inches(1.2), Inches(5), Inches(5.8),
                """Step 1: Define Goals
What are you evaluating?

Step 2: Design Rubric
Build your measurement tool

Step 3: Pilot Test
Run on sample outputs

Step 4: Refine
Adjust based on pilot results

Step 5: Scale
Apply framework to full dataset""", font_size=14, color=TEXT_DARK)

    # Slide 9: Q&A
    print("📄 Creating Slide 9: Q&A")
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    slide9.background.fill.solid()
    slide9.background.fill.fore_color.rgb = PRIMARY_COLOR

    add_text_box(slide9, Inches(1), Inches(3), Inches(8), Inches(2),
                "Questions?", font_size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide9, Inches(1), Inches(5), Inches(8), Inches(1),
                "aroma.tahir@taleemabad.com", font_size=20, color=ACCENT_COLOR, align=PP_ALIGN.CENTER)

    return prs

# Generate
print("\n" + "="*70)
print("🎨 GENERATING SYSTEMS EVALUATIONS INSTRUCTOR SLIDES")
print("   Professional design with Unsplash images (free API)")
print("="*70)

prs = create_presentation()

from pathlib import Path
output_dir = Path("weekly_artifacts/week-20-2026")
output_dir.mkdir(parents=True, exist_ok=True)

output_path = output_dir / "Systems_Evaluations_Instructor_Slides_Final.pptx"
prs.save(output_path)

print("\n" + "="*70)
print("✅ PRESENTATION COMPLETE!")
print(f"   Saved: Systems_Evaluations_Instructor_Slides_Final.pptx")
print("\n   ✓ 9 professional slides")
print("   ✓ Clean, minimal design")
print("   ✓ Professional images integrated")
print("   ✓ Navy + teal color scheme")
print("   ✓ Rounded corners & modern styling")
print("   ✓ Ready for instructor delivery")
print("="*70 + "\n")

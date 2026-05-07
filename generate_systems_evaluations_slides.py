"""
Generate Systems Evaluations presentation in PowerPoint format with images
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import requests
from PIL import Image
from io import BytesIO

# Color scheme
PRIMARY_COLOR = RGBColor(31, 78, 121)  # Dark blue
ACCENT_COLOR = RGBColor(192, 0, 0)    # Red
LIGHT_BG = RGBColor(242, 242, 242)    # Light gray
TEXT_COLOR = RGBColor(50, 50, 50)     # Dark gray
WHITE = RGBColor(255, 255, 255)

# Image URLs from Unsplash
IMAGES = {
    1: "https://images.unsplash.com/photo-1516321318423-f06a8f0d1366?w=1200&h=800&fit=crop",   # Code/Testing
    2: "https://images.unsplash.com/photo-1552664730-d307ca884978?w=1200&h=800&fit=crop",   # Problem
    3: "https://images.unsplash.com/photo-1551288049-bebda4e38f71?w=1200&h=800&fit=crop",   # Solution
    4: "https://images.unsplash.com/photo-1516321318423-f06a8f0d1366?w=1200&h=800&fit=crop",   # Metrics
    5: "https://images.unsplash.com/photo-1518932506881-b72b27e84530?w=1200&h=800&fit=crop",   # Dashboard
    6: "https://images.unsplash.com/photo-1514432324607-2e4c00038514?w=1200&h=800&fit=crop",   # Quality
    7: "https://images.unsplash.com/photo-1541961017774-22349e4a1262?w=1200&h=800&fit=crop",   # Demo
    8: "https://images.unsplash.com/photo-1516321318423-f06a8f0d1366?w=1200&h=800&fit=crop",   # Code
    9: "https://images.unsplash.com/photo-1517694712202-14dd9538aa97?w=1200&h=800&fit=crop",   # Testing
    10: "https://images.unsplash.com/photo-1519389950473-47ba0277781c?w=1200&h=800&fit=crop",  # Connection
    11: "https://images.unsplash.com/photo-1633356122544-f134324ef6db?w=1200&h=800&fit=crop",  # Team
    12: "https://images.unsplash.com/photo-1488190211105-8342881b725d?w=1200&h=800&fit=crop",  # Documentation
    13: "https://images.unsplash.com/photo-1506794778202-cad84cf45f1d?w=1200&h=800&fit=crop",  # Analysis
    14: "https://images.unsplash.com/photo-1516321318423-f06a8f0d1366?w=1200&h=800&fit=crop",  # Coding
    15: "https://images.unsplash.com/photo-1516321318423-f06a8f0d1366?w=1200&h=800&fit=crop",  # Code
    16: "https://images.unsplash.com/photo-1516321318423-f06a8f0d1366?w=1200&h=800&fit=crop",  # Code
    17: "https://images.unsplash.com/photo-1517694712202-14dd9538aa97?w=1200&h=800&fit=crop",  # Testing
    18: "https://images.unsplash.com/photo-1516321318423-f06a8f0d1366?w=1200&h=800&fit=crop",  # Analysis
    19: "https://images.unsplash.com/photo-1551288049-bebda4e38f71?w=1200&h=800&fit=crop",  # Results
    20: "https://images.unsplash.com/photo-1517694712202-14dd9538aa97?w=1200&h=800&fit=crop",  # Task
    21: "https://images.unsplash.com/photo-1516321318423-f06a8f0d1366?w=1200&h=800&fit=crop",  # Framework
    22: "https://images.unsplash.com/photo-1517694712202-14dd9538aa97?w=1200&h=800&fit=crop",  # Improvement
    23: "https://images.unsplash.com/photo-1514432324607-2e4c00038514?w=1200&h=800&fit=crop",  # Discussion
    24: "https://images.unsplash.com/photo-1551288049-bebda4e38f71?w=1200&h=800&fit=crop",  # Success
    25: "https://images.unsplash.com/photo-1633356122544-f134324ef6db?w=1200&h=800&fit=crop",  # Orchestrator
    26: "https://images.unsplash.com/photo-1495521821757-a1efb6729352?w=1200&h=800&fit=crop",  # Takeaways
    27: "https://images.unsplash.com/photo-1553288049-bebda4e38f71?w=1200&h=800&fit=crop",  # Resources
    28: "https://images.unsplash.com/photo-1543269565-cbf427effbad?w=1200&h=800&fit=crop",  # Questions
    29: "https://images.unsplash.com/photo-1516321318423-f06a8f0d1366?w=1200&h=800&fit=crop",  # Thank you
}

def add_title_slide(prs, title, subtitle, image_url):
    """Add a title slide with image background"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add image background
    try:
        img_response = requests.get(image_url, timeout=5)
        if img_response.status_code == 200:
            img = Image.open(BytesIO(img_response.content))
            img_bytes = BytesIO()
            img.save(img_bytes, format='PNG')
            img_bytes.seek(0)

            left = top = Inches(0)
            slide.shapes.add_picture(img_bytes, left, top, width=prs.slide_width, height=prs.slide_height)
    except:
        pass

    # Add semi-transparent overlay
    background = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(0, 0, 0)
    background.fill.transparency = 0.4
    background.line.color.rgb = RGBColor(0, 0, 0)

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = WHITE
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content, image_url):
    """Add a content slide with title, content, and image"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add background color
    background = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = WHITE
    background.line.color.rgb = WHITE

    # Add header bar
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.8))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.color.rgb = PRIMARY_COLOR

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(7), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE

    # Add image on the right
    try:
        img_response = requests.get(image_url, timeout=5)
        if img_response.status_code == 200:
            img = Image.open(BytesIO(img_response.content))
            img_bytes = BytesIO()
            img.save(img_bytes, format='PNG')
            img_bytes.seek(0)

            slide.shapes.add_picture(img_bytes, Inches(5.5), Inches(1.2), height=Inches(5))
    except:
        pass

    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.8), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    text_frame.text = content

    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = TEXT_COLOR
        paragraph.space_after = Pt(8)
        paragraph.level = 0

def create_presentation():
    """Create the entire Systems Evaluations presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs,
        "Systems Evaluations",
        "Knowing When Your Agent is Good\n\nAgentic AI Course | Week 20",
        IMAGES[1])

    # Slide 2: The Problem
    add_content_slide(prs,
        "The Problem",
        "You've built an agent.\nClaude can reason and act.\n\nBut: How do you know it's good?\n\n❌ No clear standard\n❌ Hard to measure quality\n❌ Can't compare agents\n❌ No early warning for problems",
        IMAGES[2])

    # Slide 3: The Solution
    add_content_slide(prs,
        "The Solution: Evaluation",
        "Systematic measurement of agent quality.\n\nThree types:\n1. Metrics (numbers)\n2. Automated checks (pass/fail)\n3. Human evaluation (expert judgment)\n\nRun continuously.\nUse results to improve.",
        IMAGES[3])

    # Slide 4: What is Evaluation?
    add_content_slide(prs,
        "What is Evaluation?",
        "Quality control for agents.\n\nYou ask:\n• Does it work?\n• Does it meet standards?\n• Where does it fail?\n• How do you improve it?\n\nEvaluation answers these.\n\nUnlike testing: evaluation measures goals achieved, not just code working.",
        IMAGES[4])

    # Slide 5: Metrics vs Checks vs Human
    add_content_slide(prs,
        "Three Types of Evaluation",
        "Metrics (quantitative):\n  Numbers: accuracy, latency, cost\n\nAutomated Checks (binary):\n  Pass/fail: glossary has 5+ terms?\n\nHuman Evaluation (qualitative):\n  Expert judgment: clarity? Rate 1-5\n\nUse all three together.",
        IMAGES[5])

    # Slide 6: The MEASURE Loop
    add_content_slide(prs,
        "The MEASURE Framework",
        "M: Measure what matters\nE: Explore failures\nA: Assess root causes\nS: Sum up health\nU: Unblock improvements\nR: Repeat\n\nThis loop never stops.\nContinuous evaluation drives improvement.",
        IMAGES[6])

    # Slide 7: Demo: Evaluating Output
    add_content_slide(prs,
        "Demo Time",
        "Watch what happens:\n\n1. I'll show an agent output\n2. Ask: Is this good? How do I know?\n3. Walk through evaluation\n4. Show metrics and rubric\n5. Make a decision\n\nDon't worry about code yet.\nJust watch the process.",
        IMAGES[7])

    # Slide 8: Metric Example
    add_content_slide(prs,
        "Example: Pass Rate Metric",
        "Definition:\n  pass_rate = (tests passed) / (total tests)\n\nExample:\n  18 tests passed, 2 failed\n  pass_rate = 18/20 = 90%\n\nWhy it matters:\n  Tracks agent reliability\n  Watch for regressions",
        IMAGES[8])

    # Slide 9: Automated Check Example
    add_content_slide(prs,
        "Example: Automated Check",
        "Rule: Glossary must have 5+ terms\n\nLogic:\nif len(glossary) < 5:\n    fail\nelse:\n    pass\n\nWhy it matters:\n  Gate before publishing\n  Prevent incomplete content",
        IMAGES[9])

    # Slide 10: Human Eval Rubric
    add_content_slide(prs,
        "Example: Human Evaluation Rubric",
        "Dimension: Clarity\n\npoor (1): Hard to understand\nfair (2): Understandable, some confusion\ngood (3): Clear, easy to follow\nexcellent (4): Very clear, well-structured\n\nEvaluator rates each output.\nAverage score = quality metric.",
        IMAGES[10])

    # Slide 11: Test Suite
    add_content_slide(prs,
        "Building a Test Suite",
        "Create test cases:\n\nTest 1: Glossary generation\n  Input: recording transcript\n  Expected: 5+ terms with definitions\n\nTest 2: Assignment creation\n  Input: learning outcomes\n  Expected: rubric with 3+ dimensions\n\nRun all tests, collect results.",
        IMAGES[11])

    # Slide 12: Drawing Room Example
    add_content_slide(prs,
        "Real-World: Drawing Room Evaluation",
        "Orchestrator metrics:\n• % sessions pass QA gate\n• Avg glossary size\n• Assignment rubric completeness\n• Publishing error rate\n\nChecks:\n• Output files exist\n• No PII in content\n• Metadata valid\n\nHuman: Review flagged items weekly",
        IMAGES[12])

    # Slide 13: Interpreting Results
    add_content_slide(prs,
        "What Do Metrics Tell You?",
        "Good news:\n  Metrics dropped -> problem signal\n\nBad news:\n  Metrics can mislead\n  Example: 95% accuracy but users unhappy\n\nAlways validate metrics against:\n  Human judgment\n  User feedback\n  Real-world impact",
        IMAGES[13])

    # Slide 14-19: Live Build sequence
    slide_titles = [
        "Live Build: Evaluator Setup",
        "Step 1: Define Metrics",
        "Step 2: Write Checks",
        "Step 3: Build Test Suite",
        "Step 4: Run Evaluation",
        "Step 5: Interpret Results"
    ]

    slide_contents = [
        "Create evaluation class.\nLoad agent output.\nPrepare test data.\n\nWe'll build this together.",
        "Pick 3 metrics:\n1. pass_rate\n2. avg_response_time\n3. content_completeness\n\nDefine calculations for each.",
        "Write pass/fail rules:\n1. Output not empty\n2. No forbidden words\n3. Matches expected schema\n\nEach rule = safety check",
        "Create TestCase objects.\nDefine inputs + expected outputs.\nRun each test.\nCollect results.",
        "Calculate metrics from results.\nCompare against baselines.\nIdentify failures.\n\nAutomated process.",
        "Read metrics.\nCompare vs threshold.\nMake decision:\n  Pass? Deploy.\n  Fail? Investigate."
    ]

    for i, (title, content) in enumerate(zip(slide_titles, slide_contents), start=14):
        add_content_slide(prs, title, content, IMAGES[i])

    # Slide 20: Assignment
    add_content_slide(prs,
        "Your Turn: Design Evaluation",
        "Scenario: Explanation Simplifier Agent\n\nYour job:\n1. Identify 3 metrics\n2. Write 3 automated checks\n3. Design human eval rubric\n4. Set acceptance criteria\n\nYou have 15 minutes.\nInstructor is here to help.",
        IMAGES[20])

    # Slide 21: Assignment Tips
    add_content_slide(prs,
        "Tips for Design Success",
        "1. Start with human judgment\n   Ask: What does good look like?\n\n2. Derive metrics from that\n   What would good outputs share?\n\n3. Write checks for obvious failures\n   What must never happen?\n\n4. Set realistic thresholds\n   When is 'good enough' good enough?",
        IMAGES[21])

    # Slide 22: Continuous Improvement
    add_content_slide(prs,
        "The Improvement Loop",
        "Week 1: Establish baseline metrics\nWeek 2: Improve prompts -> metrics go up\nWeek 3: New edge case found -> metrics drop\nWeek 4: Fix edge case -> metrics recover\n\nThis is healthy iteration.\nEvaluation drives improvement.",
        IMAGES[22])

    # Slide 23: Debrief
    add_content_slide(prs,
        "Debrief Discussion",
        "Think about:\n\n1. How would you collect data for your metrics?\n2. What's the hardest thing to measure?\n3. What happens if a metric disagrees with human judgment?",
        IMAGES[23])

    # Slide 24: Why Evaluation Matters
    add_content_slide(prs,
        "Why Evaluation Matters",
        "Without evaluation:\n  Flying blind, shipping broken agents\n\nWith evaluation:\n  Early problem detection\n  Data-driven improvement\n  Confidence in quality\n  Continuous iteration\n\nEvaluation = how you ship good agents.",
        IMAGES[24])

    # Slide 25: Drawing Room Workflow
    add_content_slide(prs,
        "Drawing Room: Evaluation in Action",
        "1. Orchestrator processes session\n2. Evaluation suite runs\n3. Metrics calculated\n4. Checks verified\n5. Human review gate\n6. Decision: keep/rebuild/kill\n7. Feedback loop improves next cycle",
        IMAGES[25])

    # Slide 26: Key Takeaways
    add_content_slide(prs,
        "Key Takeaways",
        "Evaluation measures agent quality\nUse metrics + checks + human judgment\nRun evaluation continuously\nUse results to improve\nEvaluation is how you ship good agents\n\nYou can evaluate your agents.\nYou just learned how.",
        IMAGES[26])

    # Slide 27: Resources
    add_content_slide(prs,
        "Resources & Next Steps",
        "Read:\n  Anthropic evaluation guide\n  Industry best practices\n\nTry:\n  Build evaluation for your agent\n  Measure quality metrics\n  Act on results\n\nNext: Deploy with confidence",
        IMAGES[27])

    # Slide 28: Questions
    add_content_slide(prs,
        "Questions?",
        "Ask anything:\n\n• About evaluation?\n• About your assignment?\n• About Drawing Room?\n• About metrics?\n\nNo question is too small.",
        IMAGES[28])

    # Slide 29: Thank You
    add_title_slide(prs,
        "You Just Learned Evaluation",
        "You came in asking 'How do I know if my agent is good?'\nYou're leaving knowing:\n✓ How to measure quality\n✓ How to design evaluation\n✓ How to improve systematically\n\nYou're ready to ship good agents.\nSee you next week!",
        IMAGES[29])

    # Save presentation
    output_path = r"c:\Users\Aroma Tahir\Downloads\Content Queen\weekly_artifacts\week-20-2026\Systems_Evaluations_Instructor_Slides.pptx"
    prs.save(output_path)
    return output_path

if __name__ == "__main__":
    output = create_presentation()
    print("PowerPoint created successfully: " + output)
    print("29 slides with engaging images")
    print("Professional color scheme and layout")
    print("Content optimized for 90-minute session")

"""
Systems Evaluations Instructor Slides - REDESIGNED
Best Practices Applied:
✓ Clear learning objectives
✓ Progressive complexity (simple→complex)
✓ Real examples & case studies
✓ Visual hierarchy & whitespace
✓ One key idea per slide
✓ Storytelling & narrative flow
✓ Actionable takeaways
✓ Practice scenarios
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Professional color palette
PRIMARY = RGBColor(31, 78, 121)      # Dark blue
ACCENT = RGBColor(0, 176, 176)       # Teal
SUCCESS = RGBColor(76, 175, 80)      # Green
WARNING = RGBColor(255, 152, 0)      # Orange
DANGER = RGBColor(244, 67, 54)       # Red
LIGHT_BG = RGBColor(245, 248, 250)   # Very light blue
TEXT_DARK = RGBColor(33, 33, 33)     # Very dark
TEXT_MED = RGBColor(97, 97, 97)      # Medium gray
WHITE = RGBColor(255, 255, 255)

def add_header(slide, title, subtitle=""):
    """Add clean header with optional subtitle"""
    # Background
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY
    header.line.color.rgb = PRIMARY

    # Title
    tf = header.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.15)

    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.2), Inches(10), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT

def add_text_box(slide, left, top, width, height, text, size=16, bold=False,
                 color=TEXT_DARK, align=PP_ALIGN.LEFT, line_spacing=1.5):
    """Add formatted text box"""
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP

    p = frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.line_spacing = line_spacing

    return box

def add_bullet_list(slide, left, top, width, height, items, size=16, color=TEXT_DARK):
    """Add bulleted list"""
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = frame.paragraphs[0]
        else:
            p = frame.add_paragraph()

        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.level = 0
        p.line_spacing = 1.5
        p.space_before = Pt(6)

    return box

def add_highlight_box(slide, left, top, width, height, title, content, bg_color=LIGHT_BG, accent_color=ACCENT):
    """Add highlighted content box"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.color.rgb = accent_color
    box.line.width = Pt(3)

def add_key_point(slide, left, top, text, bg_color=ACCENT):
    """Add key point callout"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(9), Inches(0.7))
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.color.rgb = bg_color

    add_text_box(slide, left + Inches(0.3), top + Inches(0.1), Inches(8.4), Inches(0.5),
                text, size=18, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

def add_example_box(slide, left, top, width, title, example_text, bg_color=RGBColor(230, 245, 250)):
    """Add example box"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(1.8))
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.color.rgb = ACCENT
    box.line.width = Pt(2)

    add_text_box(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.4),
                title, size=13, bold=True, color=ACCENT)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), Inches(1.2),
                example_text, size=11, color=TEXT_DARK)

def create_presentation():
    """Create the improved presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ==================== SLIDE 1: TITLE & LEARNING OBJECTIVES ====================
    print("📄 Slide 1: Title & Learning Objectives")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = PRIMARY

    add_text_box(s, Inches(0.5), Inches(1.5), Inches(9), Inches(1.2),
                "Systems Evaluations", size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.5), Inches(2.8), Inches(9), Inches(0.5),
                "Quality Control for AI Agents", size=28, color=ACCENT, align=PP_ALIGN.CENTER)

    add_text_box(s, Inches(1), Inches(4), Inches(8), Inches(2.8),
                "By the end of this session, you will:\n\n"
                "✓ Understand why evaluation differs from testing\n"
                "✓ Design and implement evaluation rubrics\n"
                "✓ Run hybrid (automated + manual) evaluations\n"
                "✓ Document decisions for reproducibility",
                size=18, bold=False, color=WHITE, line_spacing=1.6)

    # ==================== SLIDE 2: THE PROBLEM ====================
    print("📄 Slide 2: The Problem (Why Evaluation Matters)")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    add_header(s, "The Problem")

    add_text_box(s, Inches(0.5), Inches(1.5), Inches(9), Inches(0.8),
                "You built an AI agent. It passes all tests. But in production...",
                size=20, bold=True, color=DANGER)

    # Real scenario
    add_example_box(s, Inches(0.5), Inches(2.5), Inches(4.5),
                   "What Happened",
                   "✗ Agent hallucinates facts\n✗ Responses unhelpful\n✗ Users complain\n✗ No way to measure quality")

    add_example_box(s, Inches(5.2), Inches(2.5), Inches(4.3),
                   "Why Testing Failed",
                   "✓ Code works\n✓ Edge cases pass\n✓ Performance is fast\n✗ But user experience?")

    add_key_point(s, Inches(0.5), Inches(4.8),
                 "Tests verify code. Evaluations verify real-world value.")

    add_text_box(s, Inches(0.5), Inches(5.8), Inches(9), Inches(1.3),
                "Testing: Does the code run?\nEvaluation: Does it solve the user's problem?",
                size=16, color=TEXT_MED, align=PP_ALIGN.CENTER)

    # ==================== SLIDE 3: WHAT IS EVALUATION ====================
    print("📄 Slide 3: What is Evaluation?")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    add_header(s, "What is Evaluation?")

    add_text_box(s, Inches(0.5), Inches(1.5), Inches(9), Inches(0.8),
                "Systematic measurement of whether an agent achieves its goals",
                size=22, bold=True, color=ACCENT)

    # Three questions
    questions = [
        ("Does it work?", "Produces output without errors"),
        ("Meets standards?", "Output meets quality bar"),
        ("Where fails?", "Which use cases break?"),
        ("How improve?", "What specific changes help?")
    ]

    y_pos = 2.5
    for q, answer in questions:
        add_text_box(s, Inches(1), Inches(y_pos), Inches(0.4), Inches(0.4),
                    "→", size=20, bold=True, color=ACCENT)
        add_text_box(s, Inches(1.5), Inches(y_pos), Inches(1.5), Inches(0.4),
                    q, size=15, bold=True, color=PRIMARY)
        add_text_box(s, Inches(3.2), Inches(y_pos), Inches(6), Inches(0.4),
                    answer, size=15, color=TEXT_DARK)
        y_pos += 0.55

    # ==================== SLIDE 4: TYPES OF EVALUATION ====================
    print("📄 Slide 4: Types of Evaluation")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    add_header(s, "Three Approaches")

    # Automated
    add_text_box(s, Inches(0.5), Inches(1.5), Inches(2.8), Inches(0.5),
                "AUTOMATED", size=14, bold=True, color=SUCCESS, align=PP_ALIGN.CENTER)
    add_bullet_list(s, Inches(0.5), Inches(2.1), Inches(2.8), Inches(4.5),
                   [
                       "Fast: Seconds",
                       "Scalable: 1000s",
                       "Repeatable",
                       "Cheap",
                       "Limited nuance"
                   ], size=13, color=TEXT_DARK)

    # Manual
    add_text_box(s, Inches(3.6), Inches(1.5), Inches(2.8), Inches(0.5),
                "MANUAL", size=14, bold=True, color=WARNING, align=PP_ALIGN.CENTER)
    add_bullet_list(s, Inches(3.6), Inches(2.1), Inches(2.8), Inches(4.5),
                   [
                       "Nuanced",
                       "Context-aware",
                       "Catches edge cases",
                       "Slow: Hours",
                       "Expensive"
                   ], size=13, color=TEXT_DARK)

    # Hybrid
    add_text_box(s, Inches(6.7), Inches(1.5), Inches(2.8), Inches(0.5),
                "HYBRID ⭐", size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_bullet_list(s, Inches(6.7), Inches(2.1), Inches(2.8), Inches(4.5),
                   [
                       "Automated: First pass",
                       "Manual: Spot-check",
                       "Cost-effective",
                       "Catches both issues",
                       "Best approach"
                   ], size=13, color=TEXT_DARK)

    add_key_point(s, Inches(0.5), Inches(6.8),
                 "Industry standard: Run automated, then manual review on sample")

    # ==================== SLIDE 5: FRAMEWORK DEEP DIVE ====================
    print("📄 Slide 5: Building an Evaluation Framework")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    add_header(s, "The 5-Step Framework")

    steps = [
        ("1. DEFINE", "What goals matter?", "Quality, helpfulness, accuracy, speed"),
        ("2. DESIGN", "How measure?", "Rubrics, metrics, scoring system"),
        ("3. BUILD", "Create rubric", "Criteria → Levels → Descriptors"),
        ("4. TEST", "Pilot run", "Test on 50-100 samples first"),
        ("5. SCALE", "Full eval", "Apply to complete dataset")
    ]

    y_pos = 1.5
    for step, question, detail in steps:
        add_text_box(s, Inches(0.5), Inches(y_pos), Inches(1.2), Inches(0.35),
                    step, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Background circle for step
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(y_pos - 0.05), Inches(1.2), Inches(0.45))
        circle.fill.solid()
        circle.fill.fore_color.rgb = PRIMARY
        circle.line.color.rgb = PRIMARY
        # Re-add text on top
        tf = circle.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        add_text_box(s, Inches(1.9), Inches(y_pos), Inches(7.6), Inches(0.15),
                    question, size=12, bold=True, color=PRIMARY)
        add_text_box(s, Inches(1.9), Inches(y_pos + 0.2), Inches(7.6), Inches(0.15),
                    detail, size=11, color=TEXT_MED)

        y_pos += 0.95

    # ==================== SLIDE 6: RUBRIC EXAMPLE ====================
    print("📄 Slide 6: Rubric Example (Real Case)")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    add_header(s, "Real Rubric Example: Fact Accuracy")

    add_text_box(s, Inches(0.5), Inches(1.5), Inches(9), Inches(0.4),
                "Scenario: Evaluating if agent gives correct information",
                size=14, bold=True, color=ACCENT)

    # Rubric table as text
    rubric_data = [
        ("Level 4", "EXEMPLARY", "All facts correct, proper citations, complete answer"),
        ("Level 3", "PROFICIENT", "No major errors, 90%+ correct, mostly complete"),
        ("Level 2", "DEVELOPING", "Some errors (1-2), partial answer, confusing"),
        ("Level 1", "NEEDS WORK", "Major errors or hallucinations, incomplete")
    ]

    y_pos = 2.1
    for level, label, description in rubric_data:
        if level == "Level 4":
            color = SUCCESS
        elif level == "Level 3":
            color = WARNING
        elif level == "Level 2":
            color = RGBColor(255, 193, 7)
        else:
            color = DANGER

        add_text_box(s, Inches(0.6), Inches(y_pos), Inches(1.2), Inches(0.35),
                    label, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Colored background
        rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(y_pos), Inches(1.2), Inches(0.35))
        rect.fill.solid()
        rect.fill.fore_color.rgb = color
        rect.line.color.rgb = color
        # Re-add text
        add_text_box(s, Inches(0.6), Inches(y_pos), Inches(1.2), Inches(0.35),
                    label, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        add_text_box(s, Inches(2), Inches(y_pos), Inches(0.8), Inches(0.35),
                    level, size=10, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
        add_text_box(s, Inches(3), Inches(y_pos), Inches(6.5), Inches(0.35),
                    description, size=10, color=TEXT_DARK)

        y_pos += 0.6

    add_key_point(s, Inches(0.5), Inches(6.7),
                 "Notice: Specific descriptors → No ambiguity → Consistent scoring")

    # ==================== SLIDE 7: HOW TO AVOID BIAS ====================
    print("📄 Slide 7: Avoiding Common Mistakes")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    add_header(s, "Pitfalls & How to Avoid Them")

    pitfalls = [
        ("❌ Vague criteria", "✓ Use measurable terms", "'good' → 'factually accurate'"),
        ("❌ Rater bias", "✓ Structured rubric", "Objective scoring, not gut feeling"),
        ("❌ Too narrow sample", "✓ Representative data", "Test on diverse cases first"),
        ("❌ No documentation", "✓ Log every decision", "Why did you score it 3/4?"),
        ("❌ Single rater", "✓ Multiple judges", "Compare results, calibrate")
    ]

    y_pos = 1.5
    for bad, good, detail in pitfalls:
        add_text_box(s, Inches(0.6), Inches(y_pos), Inches(1.8), Inches(0.35),
                    bad, size=11, bold=True, color=DANGER)
        add_text_box(s, Inches(2.6), Inches(y_pos), Inches(2), Inches(0.35),
                    good, size=11, bold=True, color=SUCCESS)
        add_text_box(s, Inches(4.8), Inches(y_pos), Inches(4.7), Inches(0.35),
                    detail, size=10, color=TEXT_DARK)
        y_pos += 0.65

    # ==================== SLIDE 8: HANDS-ON EXERCISE ====================
    print("📄 Slide 8: Your Turn - Design a Rubric")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    add_header(s, "Exercise: Design Your Rubric")

    add_text_box(s, Inches(0.5), Inches(1.5), Inches(9), Inches(0.5),
                "Scenario: Evaluating customer support responses",
                size=16, bold=True, color=ACCENT)

    add_text_box(s, Inches(0.5), Inches(2.2), Inches(9), Inches(0.8),
                "What makes a GOOD support response? (Next 5 minutes: Discuss with partner)",
                size=14, color=TEXT_DARK, line_spacing=1.6)

    add_bullet_list(s, Inches(1), Inches(3.2), Inches(8), Inches(3.5),
                   [
                       "What's the main goal? (Speed? Accuracy? Empathy?)",
                       "How would you measure each criterion?",
                       "What's exemplary vs. needs work?",
                       "What would bias your judgment?",
                       "How would you document your decision?"
                   ], size=13, color=TEXT_DARK)

    # ==================== SLIDE 9: REAL RESULTS ====================
    print("📄 Slide 9: Real Results (What You Can Expect)")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    add_header(s, "When Done Right: Real Impact")

    add_example_box(s, Inches(0.5), Inches(1.6), Inches(4.4),
                   "Before Evaluation",
                   "✗ Agent seems fine\n✗ Users complain\n✗ 40% satisfaction\n✗ No visibility")

    add_example_box(s, Inches(5.3), Inches(1.6), Inches(4.2),
                   "After Evaluation",
                   "✓ Identified 3 bugs\n✓ Improved quality\n✓ 87% satisfaction\n✓ Measurable progress")

    add_key_point(s, Inches(0.5), Inches(4.2),
                 "Structured evaluation → Visible improvements → Data-driven decisions")

    add_text_box(s, Inches(0.5), Inches(5.2), Inches(9), Inches(2),
                "Example: Query about product pricing\n\n"
                "Before: \"Response is unhelpful\" (Vague)\n"
                "After: \"Misses 2 of 5 required price points\" (Measurable) → Fix it!",
                size=12, color=TEXT_DARK, line_spacing=1.5)

    # ==================== SLIDE 10: CHECKLIST ====================
    print("📄 Slide 10: Your Implementation Checklist")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    add_header(s, "Checklist: Ready to Evaluate?")

    checklist = [
        "□ Define what 'good' means for your agent (1-2 goals)",
        "□ Design 1 rubric with 4 levels + descriptors",
        "□ Pick 20-30 representative test cases",
        "□ Have 2 people score the same 5 samples",
        "□ Compare results - do they match? (Calibrate if not)",
        "□ Document scoring decisions + rationale",
        "□ Run on full test set",
        "□ Analyze results → What patterns?",
        "□ Plan improvements based on findings"
    ]

    y_pos = 1.6
    for item in checklist:
        add_text_box(s, Inches(1), Inches(y_pos), Inches(8), Inches(0.35),
                    item, size=13, color=TEXT_DARK)
        y_pos += 0.5

    add_key_point(s, Inches(0.5), Inches(6.8),
                 "Start simple. You can always refine.")

    # ==================== SLIDE 11: SUMMARY ====================
    print("📄 Slide 11: Key Takeaways")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = PRIMARY

    add_text_box(s, Inches(0.5), Inches(0.8), Inches(9), Inches(0.6),
                "What You've Learned", size=36, bold=True, color=WHITE)

    takeaways = [
        "Evaluation ≠ Testing — Measure real-world value, not just code correctness",
        "Structured rubrics eliminate bias and ensure consistency",
        "Hybrid approach works: Automated screening + Manual verification",
        "Start small: Pilot on 20-30 samples before scaling to thousands",
        "Document everything: Every decision is part of your quality audit trail"
    ]

    y_pos = 1.8
    for i, takeaway in enumerate(takeaways, 1):
        add_text_box(s, Inches(0.7), Inches(y_pos), Inches(8.6), Inches(0.8),
                    f"{i}. {takeaway}", size=14, bold=False, color=WHITE, line_spacing=1.4)
        y_pos += 1.0

    # ==================== SLIDE 12: CLOSING ====================
    print("📄 Slide 12: Questions & Next Steps")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = PRIMARY

    add_text_box(s, Inches(0.5), Inches(2), Inches(9), Inches(1.2),
                "Questions?", size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(s, Inches(0.5), Inches(3.5), Inches(9), Inches(2.5),
                "Next: Pick one of your agents and design an evaluation\n\n"
                "Send results to aroma.tahir@taleemabad.com\n\n"
                "We'll review and refine together",
                size=18, color=WHITE, align=PP_ALIGN.CENTER, line_spacing=1.6)

    return prs

# Create presentation
print("\n" + "="*70)
print("🎓 CREATING IMPROVED EDUCATIONAL SLIDES")
print("   Best practices: Clear objectives, real examples, visual hierarchy")
print("="*70 + "\n")

prs = create_presentation()

from pathlib import Path
output_dir = Path("weekly_artifacts/week-20-2026")
output_dir.mkdir(parents=True, exist_ok=True)

output_path = output_dir / "Systems_Evaluations_Instructor_Slides_Final.pptx"
prs.save(output_path)

print("="*70)
print("✅ COMPLETE! Redesigned presentation saved")
print(f"\n   File: Systems_Evaluations_Instructor_Slides_Final.pptx")
print("\n   Improvements:")
print("   ✓ Clear learning objectives")
print("   ✓ Real problem/scenario (not just theory)")
print("   ✓ Visual hierarchy & whitespace")
print("   ✓ Concrete examples & case studies")
print("   ✓ Progressive complexity")
print("   ✓ Hands-on exercise")
print("   ✓ Actionable checklist")
print("   ✓ Key takeaways highlighted")
print("="*70 + "\n")

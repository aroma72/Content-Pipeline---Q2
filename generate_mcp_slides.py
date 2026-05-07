"""
Generate MCP Servers presentation in PowerPoint format with images
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import requests
from PIL import Image
from io import BytesIO
import os

# Color scheme
PRIMARY_COLOR = RGBColor(31, 78, 121)  # Dark blue
ACCENT_COLOR = RGBColor(192, 0, 0)    # Red
LIGHT_BG = RGBColor(242, 242, 242)    # Light gray
TEXT_COLOR = RGBColor(50, 50, 50)     # Dark gray
WHITE = RGBColor(255, 255, 255)

# Image URLs from Unsplash and free sources
IMAGES = {
    1: "https://images.unsplash.com/photo-1516321318423-f06a8f0d1366?w=1200&h=800&fit=crop",  # Programming
    2: "https://images.unsplash.com/photo-1517694712202-14dd9538aa97?w=1200&h=800&fit=crop",  # Code duplication
    3: "https://images.unsplash.com/photo-1519389950473-47ba0277781c?w=1200&h=800&fit=crop",  # Connected solutions
    4: "https://images.unsplash.com/photo-1516321318423-f06a8f0d1366?w=1200&h=800&fit=crop",  # Protocol
    5: "https://images.unsplash.com/photo-1558618666-fcd25c85cd64?w=1200&h=800&fit=crop",    # Architecture
    6: "https://images.unsplash.com/photo-1495521821757-a1efb6729352?w=1200&h=800&fit=crop",  # Restaurant
    7: "https://images.unsplash.com/photo-1517694712202-14dd9538aa97?w=1200&h=800&fit=crop",  # Demo
    8: "https://images.unsplash.com/photo-1484480974693-6ca0a78fb36b?w=1200&h=800&fit=crop",  # Server
    9: "https://images.unsplash.com/photo-1516321318423-f06a8f0d1366?w=1200&h=800&fit=crop",  # Code
    10: "https://images.unsplash.com/photo-1517694712202-14dd9538aa97?w=1200&h=800&fit=crop", # Connection
    11: "https://images.unsplash.com/photo-1519389950473-47ba0277781c?w=1200&h=800&fit=crop", # Loop
    12: "https://images.unsplash.com/photo-1633356122544-f134324ef6db?w=1200&h=800&fit=crop", # Orchestration
    13: "https://images.unsplash.com/photo-1517694712202-14dd9538aa97?w=1200&h=800&fit=crop", # Transport
    14: "https://images.unsplash.com/photo-1516321318423-f06a8f0d1366?w=1200&h=800&fit=crop", # Live build
    15: "https://images.unsplash.com/photo-1516321318423-f06a8f0d1366?w=1200&h=800&fit=crop", # Tool
    16: "https://images.unsplash.com/photo-1516321318423-f06a8f0d1366?w=1200&h=800&fit=crop", # Execution
    17: "https://images.unsplash.com/photo-1484480974693-6ca0a78fb36b?w=1200&h=800&fit=crop", # Server start
    18: "https://images.unsplash.com/photo-1516321318423-f06a8f0d1366?w=1200&h=800&fit=crop", # Connect
    19: "https://images.unsplash.com/photo-1517694712202-14dd9538aa97?w=1200&h=800&fit=crop", # Loop action
    20: "https://images.unsplash.com/photo-1517694712202-14dd9538aa97?w=1200&h=800&fit=crop", # Assignment
    21: "https://images.unsplash.com/photo-1516321318423-f06a8f0d1366?w=1200&h=800&fit=crop", # Tips
    22: "https://images.unsplash.com/photo-1517694712202-14dd9538aa97?w=1200&h=800&fit=crop", # Stretch
    23: "https://images.unsplash.com/photo-1514984879728-be823260b4d6?w=1200&h=800&fit=crop", # Debrief
    24: "https://images.unsplash.com/photo-1519389950473-47ba0277781c?w=1200&h=800&fit=crop", # Why matters
    25: "https://images.unsplash.com/photo-1633356122544-f134324ef6db?w=1200&h=800&fit=crop", # Drawing Room
    26: "https://images.unsplash.com/photo-1495521821757-a1efb6729352?w=1200&h=800&fit=crop", # Takeaways
    27: "https://images.unsplash.com/photo-1553288049-bebda4e38f71?w=1200&h=800&fit=crop",    # Resources
    28: "https://images.unsplash.com/photo-1543269565-cbf427effbad?w=1200&h=800&fit=crop",    # Questions
    29: "https://images.unsplash.com/photo-1516321318423-f06a8f0d1366?w=1200&h=800&fit=crop", # Thank you
}

def download_image(url):
    """Download image from URL and return PIL Image"""
    try:
        response = requests.get(url, timeout=5)
        if response.status_code == 200:
            return Image.open(BytesIO(response.content))
    except:
        pass
    return None

def add_title_slide(prs, title, subtitle, image_url):
    """Add a title slide with image background"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add image background
    try:
        img_response = requests.get(image_url, timeout=5)
        if img_response.status_code == 200:
            img = Image.open(BytesIO(img_response.content))
            img_bytes = BytesIO()
            img.save(img_bytes, format='PNG')
            img_bytes.seek(0)

            left = top = Inches(0)
            slide.shapes.add_picture(img_bytes, left, top, width=prs.slide_width, height=prs.slide_height)
    except:
        pass

    # Add semi-transparent overlay
    background = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(0, 0, 0)
    background.fill.transparency = 0.4
    background.line.color.rgb = RGBColor(0, 0, 0)

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = WHITE
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content, image_url):
    """Add a content slide with title, content, and image"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add background color
    background = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = WHITE
    background.line.color.rgb = WHITE

    # Add header bar
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.8))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.color.rgb = PRIMARY_COLOR

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(7), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE

    # Add image on the right
    try:
        img_response = requests.get(image_url, timeout=5)
        if img_response.status_code == 200:
            img = Image.open(BytesIO(img_response.content))
            img_bytes = BytesIO()
            img.save(img_bytes, format='PNG')
            img_bytes.seek(0)

            slide.shapes.add_picture(img_bytes, Inches(5.5), Inches(1.2), height=Inches(5))
    except:
        pass

    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.8), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    text_frame.text = content

    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = TEXT_COLOR
        paragraph.space_after = Pt(8)
        paragraph.level = 0

def create_presentation():
    """Create the entire MCP Servers presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs,
        "Model Context Protocol",
        "Giving AI Agents Hands\n\nAgentic AI Course | Week 19",
        IMAGES[1])

    # Slide 2: The Problem
    add_content_slide(prs,
        "The Problem: Without MCP",
        "❌ Code duplication across agents\n❌ Inconsistent implementations\n❌ Updates scattered everywhere\n\nEvery agent rewrites the same tools:\n• read_file()\n• write_file()\n• search()\n\nFinding a bug? Fix it in 3 places.",
        IMAGES[2])

    # Slide 3: The Solution
    add_content_slide(prs,
        "The Solution: With MCP",
        "✅ One source of truth\n✅ Reusable across agents\n✅ Update once, benefit everywhere\n\nWrite tools once. Connect any agent.\n\nAgent code never changes when you add tools.",
        IMAGES[3])

    # Slide 4: What is MCP
    add_content_slide(prs,
        "What is MCP?",
        "Model Context Protocol: A standard protocol for connecting external tools to Claude.\n\nKey Idea:\n• Tools live on a separate server\n• Your agent's client connects to it\n• Claude decides when to use tools\n\nLike USB for AI agents.",
        IMAGES[4])

    # Slide 5: Architecture Diagram
    add_content_slide(prs,
        "MCP Architecture",
        "Three components:\n\n1. Host (Your Application)\n   • Holds the MCP client\n\n2. MCP Client (in your agent)\n   • Talks to the server\n   • Fetches available tools\n\n3. MCP Server (separate process)\n   • Defines and executes tools\n   • Returns results",
        IMAGES[5])

    # Slide 6: Restaurant Analogy
    add_content_slide(prs,
        "The Restaurant Analogy",
        "Customer (Claude): Places an order\n↓\nWaiter (MCP Client): Takes the order\n↓\nKitchen (MCP Server): Makes the food\n↓\nWaiter (MCP Client): Brings it back\n↓\nCustomer (Claude): Enjoys the meal\n\nKey: Customer never goes to the kitchen.",
        IMAGES[6])

    # Slide 7: Demo Time Part 1
    add_content_slide(prs,
        "Demo Time",
        "Watch what happens:\n\n1. I'll run an MCP server\n2. Connect Claude via the SDK\n3. Ask Claude a question\n4. Claude autonomously uses the server\n\nDon't worry about details yet. Just watch.",
        IMAGES[7])

    # Slide 8: Demo Time Part 2
    add_content_slide(prs,
        "What Just Happened",
        "The autonomous flow:\n\n1. User asks question\n2. Claude identifies need for tool\n3. Client asks server for tools\n4. Server sends tool list\n5. Claude chooses appropriate tool\n6. Client executes on server\n7. Server returns result\n8. Claude uses result\n\n✅ All autonomous. No manual steps.",
        IMAGES[8])

    # Slide 9: MCP Server Structure
    add_content_slide(prs,
        "Building an MCP Server",
        "Two essential methods:\n\n1. list_tools()\n   → Tell Claude what tools you have\n\n2. call_tool(name, arguments)\n   → Execute the tool\n\nThat's the whole server.\n\nInside call_tool(): check the name, do the work, return result.",
        IMAGES[9])

    # Slide 10: MCP Client Connection
    add_content_slide(prs,
        "Connecting Your Agent",
        "Three steps:\n\n1. Initialize client connection\n2. Fetch tool list from server\n3. Pass tools to Claude\n\nClaude sees the tools and uses them autonomously.\n\nAgent code stays simple. Loose coupling.",
        IMAGES[10])

    # Slide 11: The Agentic Loop
    add_content_slide(prs,
        "The Agentic Loop",
        "How MCP powers autonomous agents:\n\n1. Claude thinks: What can I do?\n2. Client fetches tool list\n3. Claude decides: I'll use this tool\n4. Client calls server\n5. Server executes, returns result\n6. Claude uses result in reasoning\n7. Loop until task complete\n\nWithout MCP: Claude can think but can't act.",
        IMAGES[11])

    # Slide 12: Drawing Room Example
    add_content_slide(prs,
        "Real-World: Drawing Room",
        "Orchestrator connects to multiple MCP servers:\n\n📁 FileSystem Server\n   • read_file()\n   • write_file()\n\n📚 LMS Server\n   • get_submissions()\n   • post_assignments()\n\n🎬 Video Tools Server\n   • transcode()\n   • extract_audio()\n\nEach team manages their own server.",
        IMAGES[12])

    # Slide 13: Transport Types
    add_content_slide(prs,
        "Transport: Stdio vs SSE",
        "Two ways to communicate:\n\n📍 Stdio (Local)\n   • Server is a local process\n   • Talk via stdin/stdout\n   • Simple, no network\n\n☁️ SSE (Remote)\n   • Server runs in the cloud\n   • Talk via HTTP\n   • Good for resource-heavy tools\n\nToday: We use stdio.",
        IMAGES[13])

    # Slide 14-19: Live Build sequence
    for i in range(14, 20):
        if i == 14:
            add_content_slide(prs,
                "Live Build: Let's Code",
                "We're building a minimal MCP server together.\n\nYou'll see:\n1. Tool definition\n2. Tool execution\n3. Connection to Claude\n\nWatch first. You'll build in the assignment.",
                IMAGES[14])
        elif i == 15:
            add_content_slide(prs,
                "Step 1: Define the Tool",
                "In list_tools():\n\nname: 'get_today'\ndescription: 'Returns today's date'\ninputSchema: {empty}\n\nThis tells Claude: I have a tool called get_today.",
                IMAGES[15])
        elif i == 16:
            add_content_slide(prs,
                "Step 2: Execute the Tool",
                "In call_tool():\n\nif name == 'get_today':\n  return today's date\n\nWhen Claude calls it, this code runs.",
                IMAGES[16])
        elif i == 17:
            add_content_slide(prs,
                "Step 3: Start Server",
                "asyncio.run(stdio_server(server))\n\nServer is now running.\nListening for client requests.\nWaiting for Claude to connect.",
                IMAGES[17])
        elif i == 18:
            add_content_slide(prs,
                "Step 4: Connect Claude",
                "Client connects to server.\nPasses tools to Claude.\nAsks: 'What's today's date?'\n\nClaude autonomously:\n• Sees the tool\n• Decides to use it\n• Gets the result",
                IMAGES[18])
        elif i == 19:
            add_content_slide(prs,
                "Step 5: The Loop in Action",
                "User: 'What's today's date?'\nClaude: 'I'll use get_today'\nServer: 'Today is [date]'\nClaude: 'Today is [date]'\n\nAutonomous. No manual steps.",
                IMAGES[19])

    # Slide 20: Assignment
    add_content_slide(prs,
        "Your Turn: The Assignment",
        "Add a search_files(keyword) tool:\n\n✓ Accept keyword argument\n✓ Search files in workspace\n✓ Return matching filenames\n✓ Return 'No matches' if none\n\nYou have 15 minutes.\nInstructor is here to help.",
        IMAGES[20])

    # Slide 21: Assignment Tips
    add_content_slide(prs,
        "Assignment Tips",
        "1. Define in list_tools()\n   name, description, inputSchema\n\n2. Implement in call_tool()\n   Add if-branch for search_files\n\n3. Test it\n   Run server, ask Claude to search\n\nYou've got this. Same pattern as demo.",
        IMAGES[21])

    # Slide 22: Stretch Goal
    add_content_slide(prs,
        "Stretch Goal",
        "Add summarise_file(filename):\n\n✓ Takes filename\n✓ Reads the file\n✓ Calls Claude for 1-paragraph summary\n✓ Returns summary\n\nHint: Use anthropic.Anthropic() from within the tool.",
        IMAGES[22])

    # Slide 23: Debrief Questions
    add_content_slide(prs,
        "Debrief Discussion",
        "Think about these:\n\n1. How much agent code changed when you added the tool?\n   (Answer: Zero)\n\n2. What if two agents connected to this server?\n   (Answer: Both see same tools)\n\n3. What MCP server would you build for Drawing Room?\n   (Brainstorm...)",
        IMAGES[23])

    # Slide 24: Why MCP Matters
    add_content_slide(prs,
        "Why MCP Matters",
        "Without MCP:\n❌ Duplication, tight coupling, no sharing\n\nWith MCP:\n✅ One source of truth\n✅ Loose coupling\n✅ Independent updates\n✅ Team collaboration\n\nMCP is how you build scalable, maintainable agents.",
        IMAGES[24])

    # Slide 25: Drawing Room Workflow
    add_content_slide(prs,
        "Drawing Room: From Signals to Assets",
        "Orchestrator workflow:\n\n📖 Read signal_backlog\n📋 Plan content\n📝 Write content map\n📚 Generate learner packs\n🎬 Video processing\n📤 Publish assets\n\nMCP is the glue. It's how the orchestrator touches the real world.",
        IMAGES[25])

    # Slide 26: Key Takeaways
    add_content_slide(prs,
        "Key Takeaways",
        "✅ MCP is a protocol for tools\n✅ Servers expose, clients connect\n✅ Claude decides autonomously\n✅ You can build MCP servers\n✅ MCP scales across agents\n✅ Drawing Room uses MCP\n\nYou can build this. You just did.",
        IMAGES[26])

    # Slide 27: Resources
    add_content_slide(prs,
        "Resources & Next Steps",
        "📚 Official Docs:\n   modelcontextprotocol.io\n\n💻 Examples:\n   github.com/anthropics/mcp\n\n🏗️ Try at Home:\n   Build your own MCP server\n\n🚀 Next Session:\n   Deployment, multi-server architectures",
        IMAGES[27])

    # Slide 28: Questions
    add_content_slide(prs,
        "Questions?",
        "Ask anything:\n\n• About MCP?\n• About the assignment?\n• About Drawing Room?\n• About agents?\n\nNo question is too small.\nLet's discuss.",
        IMAGES[28])

    # Slide 29: Thank You
    add_title_slide(prs,
        "You Just Learned MCP",
        "You came in asking 'What's MCP?'\nYou're leaving knowing:\n✅ What it is • ✅ How it works • ✅ How to build\n✅ How to connect • ✅ Why it matters\n\nYou're ready to build with MCP.\nSee you next week!",
        IMAGES[29])

    # Save presentation
    output_path = r"c:\Users\Aroma Tahir\Downloads\Content Queen\weekly_artifacts\week-19-2026\MCP_Servers_Instructor_Slides.pptx"
    prs.save(output_path)
    return output_path

if __name__ == "__main__":
    output = create_presentation()
    print("PowerPoint created successfully: " + output)
    print("29 slides with engaging images")
    print("Professional color scheme and layout")
    print("Speaker notes integrated into content")

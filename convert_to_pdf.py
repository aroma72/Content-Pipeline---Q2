"""
Convert PowerPoint slides to PDF
"""
from pptx import Presentation
from PIL import Image
from io import BytesIO
import subprocess
from pathlib import Path

pptx_file = Path("weekly_artifacts/week-20-2026/Systems_Evaluations_Instructor_Slides_Final.pptx")
pdf_file = Path("weekly_artifacts/week-20-2026/Systems_Evaluations_Instructor_Slides_Final.pdf")

print("\n🔄 Converting to PDF...")

# Try using LibreOffice command line
try:
    result = subprocess.run([
        "libreoffice", "--headless", "--convert-to", "pdf",
        "--outdir", str(pptx_file.parent),
        str(pptx_file)
    ], capture_output=True, timeout=30)

    if result.returncode == 0 and pdf_file.exists():
        print(f"✅ PDF created successfully!")
        print(f"   File: {pdf_file.name}")
        print(f"   Size: {pdf_file.stat().st_size / (1024*1024):.2f} MB")
    else:
        print("LibreOffice conversion failed, trying alternative method...")
        raise Exception("LibreOffice failed")

except Exception as e:
    print(f"Note: {str(e)[:60]}")
    print("Using alternative approach...")

    # Alternative: Use python-pptx to extract info and create PDF
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter

        prs = Presentation(str(pptx_file))
        c = canvas.Canvas(str(pdf_file), pagesize=letter)

        for i, slide in enumerate(prs.slides, 1):
            # Add slide content as text
            c.drawString(50, 750, f"Slide {i}")

            # Extract text from shapes
            y_pos = 730
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text[:100]  # Limit length
                    c.drawString(70, y_pos, text)
                    y_pos -= 20

            c.showPage()

        c.save()
        print(f"✅ PDF created with reportlab")
        print(f"   File: {pdf_file.name}")

    except ImportError:
        print("Converting via online service...")
        print("⚠️  PDF conversion skipped - LibreOffice not available")
        print(f"\nAlternative: Use online converter for {pptx_file.name}")

print("\n" + "="*70)

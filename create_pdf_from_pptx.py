"""
Create PDF from PowerPoint presentation
Extracts text and formats professionally
"""

from pptx import Presentation
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.units import inch
from reportlab.lib.colors import HexColor
from pathlib import Path

# Colors
PRIMARY = HexColor("#1F4E79")      # Dark blue
ACCENT = HexColor("#00B0B0")       # Teal
TEXT_DARK = HexColor("#323232")    # Dark gray

def extract_slide_text(slide):
    """Extract all text from a slide"""
    lines = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            text = shape.text.strip()
            if text:
                lines.append(text)
    return lines

def create_pdf():
    """Create PDF from PPTX"""
    pptx_path = Path("weekly_artifacts/week-20-2026/Systems_Evaluations_Instructor_Slides_Final.pptx")
    pdf_path = Path("weekly_artifacts/week-20-2026/Systems_Evaluations_Instructor_Slides_Final.pdf")

    prs = Presentation(str(pptx_path))

    # Create canvas
    c = canvas.Canvas(str(pdf_path), pagesize=letter)
    width, height = letter

    print(f"Creating PDF with {len(prs.slides)} slides...")

    for slide_num, slide in enumerate(prs.slides, 1):
        # Header
        c.setFillColor(PRIMARY)
        c.rect(0, height - 0.5*inch, width, 0.5*inch, fill=1)

        c.setFont("Helvetica-Bold", 24)
        c.setFillColor(HexColor("#FFFFFF"))
        c.drawString(0.4*inch, height - 0.3*inch, f"Systems Evaluations — Slide {slide_num}")

        # Accent line
        c.setStrokeColor(ACCENT)
        c.setLineWidth(3)
        c.line(0, height - 0.52*inch, width, height - 0.52*inch)

        # Extract and draw text
        lines = extract_slide_text(slide)

        y_pos = height - 1*inch
        c.setFont("Helvetica", 11)
        c.setFillColor(TEXT_DARK)

        for line in lines:
            if len(line) > 100:  # Break long lines
                words = line.split()
                current_line = ""
                for word in words:
                    if len(current_line + " " + word) > 80:
                        if y_pos < 0.5*inch:  # New page
                            c.showPage()
                            y_pos = height - 0.6*inch
                            # Re-add header for new page
                            c.setFillColor(PRIMARY)
                            c.rect(0, height - 0.5*inch, width, 0.5*inch, fill=1)
                            c.setFont("Helvetica-Bold", 18)
                            c.setFillColor(HexColor("#FFFFFF"))
                            c.drawString(0.4*inch, height - 0.3*inch, f"Systems Evaluations (continued)")
                            c.setFont("Helvetica", 11)
                            c.setFillColor(TEXT_DARK)
                        c.drawString(0.5*inch, y_pos, current_line)
                        y_pos -= 0.2*inch
                        current_line = word
                    else:
                        current_line += " " + word if current_line else word

                if current_line:
                    if y_pos < 0.5*inch:
                        c.showPage()
                        y_pos = height - 0.6*inch
                    c.drawString(0.5*inch, y_pos, current_line)
                    y_pos -= 0.2*inch
            else:
                if y_pos < 0.5*inch:
                    c.showPage()
                    y_pos = height - 0.6*inch
                    # Re-add header
                    c.setFillColor(PRIMARY)
                    c.rect(0, height - 0.5*inch, width, 0.5*inch, fill=1)
                    c.setFont("Helvetica-Bold", 18)
                    c.setFillColor(HexColor("#FFFFFF"))
                    c.drawString(0.4*inch, height - 0.3*inch, f"Systems Evaluations (continued)")
                    c.setFont("Helvetica", 11)
                    c.setFillColor(TEXT_DARK)

                if line.isupper() and len(line) < 40:  # Title style
                    c.setFont("Helvetica-Bold", 13)
                    c.drawString(0.5*inch, y_pos, line)
                    c.setFont("Helvetica", 11)
                else:
                    c.drawString(0.5*inch, y_pos, line)

                y_pos -= 0.25*inch

        # Footer
        c.setFont("Helvetica", 9)
        c.setFillColor(HexColor("#999999"))
        c.drawString(0.4*inch, 0.3*inch, f"Systems Evaluations Instructor Guide | Page {slide_num}")
        c.drawString(width - 1.5*inch, 0.3*inch, f"aroma.tahir@taleemabad.com")

        if slide_num < len(prs.slides):
            c.showPage()

    c.save()

    print(f"✅ PDF created successfully!")
    print(f"   File: {pdf_path.name}")
    size_kb = pdf_path.stat().st_size / 1024
    print(f"   Size: {size_kb:.1f} KB")
    print(f"\n   Location: weekly_artifacts/week-20-2026/")

    return str(pdf_path)

if __name__ == "__main__":
    print("\n" + "="*70)
    print("📄 CONVERTING POWERPOINT TO PDF")
    print("="*70 + "\n")

    pdf_path = create_pdf()

    print("\n" + "="*70)
    print("Ready to share!")
    print("="*70 + "\n")
